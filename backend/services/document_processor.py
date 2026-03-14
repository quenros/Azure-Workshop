import os
import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation

class DocumentProcessor:
    def extract_text(self, file_stream, filename):
        """
        Extracts raw text from the uploaded file stream based on its extension.
        """
        ext = os.path.splitext(filename)[1].lower()
        text = ""
        
        try:
            # ---> THE FIX: Rewind the file pointer to the beginning <---
            file_stream.seek(0)
            
            # Read the binary data from the Flask FileStorage object
            blob_data = file_stream.read()

            if ext == '.pdf':
                with io.BytesIO(blob_data) as temp_file:
                    document = fitz.open("pdf", temp_file)
                    for page_num in range(len(document)):
                        page = document.load_page(page_num)
                        text += page.get_text()
                        
            elif ext == '.docx':
                with io.BytesIO(blob_data) as temp_file:
                    doc = docx.Document(temp_file)
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                        
            elif ext == '.pptx':
                with io.BytesIO(blob_data) as temp_file:
                    presentation = Presentation(temp_file)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                                
            elif ext in ['.txt', '.csv']:
                with io.BytesIO(blob_data) as temp_file:
                    text = temp_file.read().decode('utf-8', errors='replace')
                    
        except Exception as e:
            print(f"Error extracting text from {filename}: {str(e)}")
            
        # Reset the file pointer again just in case another service needs it later
        file_stream.seek(0)
        
        return text.strip()